"""Generate QA Automation PPT — current 'start' command workflow."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0B, 0x2C, 0x5E)
CYAN    = RGBColor(0x00, 0xAE, 0xEF)
ORANGE  = RGBColor(0xE8, 0x53, 0x2A)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
LGREY   = RGBColor(0xF4, 0xF7, 0xFB)
MGREY   = RGBColor(0xCC, 0xD6, 0xE8)
DGREY   = RGBColor(0x44, 0x55, 0x70)

SZ_TITLE    = 28
SZ_SUB      = 16   # sub-heading — as requested
SZ_BODY     = 13
SZ_SMALL    = 11
SZ_TABLE_H  = 12
SZ_TABLE_B  = 11

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    return s

def txt(slide, text, l, t, w, h, size=SZ_BODY, bold=False,
        color=DARK, align=PP_ALIGN.LEFT, italic=False, mono=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Courier New" if mono else "Calibri"
    return tb

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, fill=NAVY)
    txt(slide, title, 0.35, 0.10, 12.6, 0.75,
        size=SZ_TITLE, bold=True, color=WHITE)
    add_rect(slide, 0, 1.35, 13.33, 0.06, fill=CYAN)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.82, 12.6, 0.45,
            size=SZ_SUB, color=CYAN)
    return 1.55 if subtitle else 1.45

def sub(slide, text, l, t, w=12.5, color=NAVY):
    txt(slide, text, l, t, w, 0.40, size=SZ_SUB, bold=True, color=color)

def bullets(slide, items, l, t, w, size=SZ_BODY, color=DARK, gap=0.48):
    for i, item in enumerate(items):
        txt(slide, "▪  " + item, l, t + i * gap, w, gap - 0.04,
            size=size, color=color)

def make_table(slide, heads, rows, l, t, w, col_w=None):
    nr = len(rows) + 1; nc = len(heads)
    tbl = slide.shapes.add_table(nr, nc,
          Inches(l), Inches(t), Inches(w), Inches(0.38 * nr)).table
    tbl.first_row = True
    if col_w:
        for i, cw in enumerate(col_w): tbl.columns[i].width = Inches(cw)
    for ci, h in enumerate(heads):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        r = c.text_frame.paragraphs[0].add_run()
        r.text = h; r.font.size = Pt(SZ_TABLE_H)
        r.font.bold = True; r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        bg = RGBColor(0xEB, 0xF3, 0xFB) if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci); c.fill.solid(); c.fill.fore_color.rgb = bg
            r = c.text_frame.paragraphs[0].add_run()
            r.text = val; r.font.size = Pt(SZ_TABLE_B)
            r.font.color.rgb = DARK

def pill(slide, label, l, t, w, h, bg=NAVY, fg=WHITE, size=SZ_SMALL):
    add_rect(slide, l, t, w, h, fill=bg)
    txt(slide, label, l, t, w, h, size=size, bold=True, color=fg,
        align=PP_ALIGN.CENTER)

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)

def bg(s): add_rect(s, 0, 0, 13.33, 7.5, fill=LGREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 5.5, 13.33, 2.0, fill=CYAN)

txt(s, "Jira-Driven QA Automation",
    0.7, 1.3, 11.9, 1.2, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Automated Test Generation, Review & Execution — driven by a single  start  command",
    0.7, 2.7, 11.9, 0.6, size=SZ_SUB, color=CYAN, align=PP_ALIGN.CENTER)

txt(s, "LBLR Voice Services  ·  Project: LBVOICESER",
    0.7, 3.5, 11.9, 0.5, size=SZ_BODY, italic=True,
    color=MGREY, align=PP_ALIGN.CENTER)

txt(s, "Java 11  ·  REST Assured  ·  JUnit 5  ·  Maven  ·  Allure  ·  Jira / Xray Cloud  ·  Claude Code",
    0.7, 5.6, 11.9, 0.5, size=SZ_BODY, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview: how the start command works
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "How It Works — The  start  Command",
           "Type  start  in the Claude Code chat inside VS Code. Everything below runs automatically.")

sub(s, "Four phases — per-ticket processing, then one combined test run:", 0.35, y + 0.05)

phases = [
    (NAVY,   "Phase 1",  "JQL Fetch",
     "Pulls every LBVOICESER story in 'Ready for testing' with no qa-auto-generated or qa-context-requested label."),
    (CYAN,   "Phase 2",  "Context Check",
     "Evaluates each ticket: summary · description · ≥1 AC · HTTP method · endpoint · request/response example."),
    (GREEN,  "Phase 3",  "SUFFICIENT → Gates 1 & 2",
     "Test case table → human approval → Jira Test issue created → script written → human approval → queued for run."),
    (ORANGE, "Phase 3b", "INSUFFICIENT → Comment",
     "Posts a structured ADF comment listing exactly what's missing. Stamps qa-context-requested on the Jira ticket."),
    (NAVY,   "Phase 4",  "Gate 3 — Combined Maven Run",
     "One  mvn test  command covering all approved classes. Allure report opens. Jira auto-updated per class."),
]

ROW_H = 0.82
for i, (col, phase, title, desc) in enumerate(phases):
    ty = y + 0.52 + i * ROW_H
    pill(s, phase, 0.35, ty + 0.08, 1.20, 0.52, bg=col)
    add_rect(s, 1.60, ty, 11.35, ROW_H - 0.08, fill=WHITE, line=MGREY)
    txt(s, title, 1.72, ty + 0.05, 11.1, 0.34, size=SZ_SUB, bold=True, color=NAVY)
    txt(s, desc,  1.72, ty + 0.40, 11.1, 0.34, size=SZ_SMALL, italic=True, color=DGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — JQL & Context Check
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Phase 1 & 2 — Fetch Tickets & Context Check")

# Left: JQL
sub(s, "JQL Query (Phase 1)", 0.35, y + 0.05)
add_rect(s, 0.35, y + 0.52, 5.9, 1.55, fill=RGBColor(0x0D, 0x1A, 0x2E), line=CYAN)
jql_lines = [
    "project = LBVOICESER",
    "AND status = \"Ready for testing\"",
    "AND issuetype in (Story, Task)",
    "AND (labels is EMPTY OR",
    "  labels NOT IN (\"qa-auto-generated\",",
    "  \"qa-context-requested\"))",
]
for i, line in enumerate(jql_lines):
    txt(s, line, 0.50, y + 0.60 + i * 0.22, 5.6, 0.24,
        size=SZ_SMALL, color=GREEN, mono=True)

# Right: mandatory fields
sub(s, "Mandatory Fields (Phase 2 Context Check)", 6.55, y + 0.05)
make_table(s,
    ["Field", "Required for"],
    [
        ["Summary",              "All tickets"],
        ["Description",          "All tickets"],
        ["≥1 Acceptance Criterion", "All tickets"],
        ["HTTP Method",          "All tickets"],
        ["Endpoint Path",        "All tickets"],
        ["Request body example", "POST / PUT / PATCH"],
        ["Sample response body", "GET"],
    ],
    6.55, y + 0.52, 6.42,
    col_w=[2.8, 3.45],
)

# Bottom: outcome labels
y2 = y + 2.30
sub(s, "Outcome labels — permanent ticket markers:", 0.35, y2)
make_table(s,
    ["Label", "Meaning", "Effect on future syncs"],
    [
        ["qa-auto-generated",   "Test generated & approved",
         "Ticket excluded forever — never processed again"],
        ["qa-context-requested","Comment posted, info missing",
         "Excluded until assignee removes the label"],
    ],
    0.35, y2 + 0.48, 12.6,
    col_w=[2.9, 4.0, 5.5],
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — INSUFFICIENT Path
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Phase 3b — INSUFFICIENT Ticket: Comment & Label",
           "When a ticket lacks required fields, Claude posts a structured comment before any QA time is spent.")

sub(s, "What gets posted to Jira (ADF comment):", 0.35, y + 0.05)
make_table(s,
    ["Section in comment", "Content"],
    [
        ["Mention @assignee",  "Tags the ticket assignee directly"],
        ["Missing fields list","Exact field names that are absent (e.g. 'HTTP Method', 'Endpoint Path')"],
        ["Examples provided",  "Shows what a well-formed AC or request body looks like"],
        ["No-action note",     "Explains the ticket will be re-evaluated automatically once updated"],
    ],
    0.35, y + 0.52, 7.2,
    col_w=[2.8, 4.25],
)

sub(s, "After comment — automatic label stamp:", 7.80, y + 0.05)
add_rect(s, 7.80, y + 0.52, 5.1, 0.62, fill=ORANGE)
txt(s, "qa-context-requested", 7.90, y + 0.60, 4.9, 0.46,
    size=SZ_SUB, bold=True, color=WHITE)

bullets(s, [
    "JQL filter excludes this label on every subsequent sync",
    "No duplicate comments — ticket is processed exactly once",
    "Automatically re-enters the queue once assignee fixes the",
    "  ticket and removes the label",
], 7.80, y + 1.25, 5.1, size=SZ_SMALL, color=DARK, gap=0.46)

# Flow diagram
y3 = y + 2.65
sub(s, "Ticket lifecycle — INSUFFICIENT path:", 0.35, y3)
steps_i = ["Ticket in 'Ready for testing'", "Fails context check", "ADF comment posted",
           "qa-context-requested label", "Excluded from next sync", "Re-enters when fixed"]
w_s = 1.95
for i, st in enumerate(steps_i):
    col = NAVY if i % 2 == 0 else ORANGE
    add_rect(s, 0.35 + i * (w_s + 0.18), y3 + 0.48, w_s, 0.62, fill=col)
    txt(s, st, 0.35 + i * (w_s + 0.18), y3 + 0.52, w_s, 0.54,
        size=SZ_SMALL, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps_i) - 1:
        txt(s, "→", 0.35 + (i + 1) * (w_s + 0.18) - 0.22, y3 + 0.60, 0.22, 0.38,
            size=14, bold=True, color=DGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Gate 1: Test Case Review
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Gate 1 — Test Case Review  (Human Approval)",
           "Claude presents a test case table first, then asks for approval — no code written yet.")

# Left column: what Claude shows
sub(s, "Step 1 — Claude shows the test case table:", 0.35, y + 0.05)
make_table(s,
    ["Test method name", "DisplayName", "AC it covers"],
    [
        ["missingAuthHeaderReturns400", "No auth header returns 400", "AC #1"],
        ["emptyBodyReturns400",         "Empty body returns 400",     "AC #2"],
        ["missingVoiceFieldReturns400", "Missing voice field returns 400", "AC #3"],
        ["missingSpeechTextReturns400", "Missing speechText returns 400",  "AC #4"],
    ],
    0.35, y + 0.52, 6.5,
    col_w=[2.6, 2.5, 1.25],
)

sub(s, "Step 2 — AskUserQuestion (UI dropdown):", 0.35, y + 2.40)
add_rect(s, 0.35, y + 2.85, 6.5, 0.62, fill=RGBColor(0xEB, 0xF3, 0xFB), line=CYAN)
txt(s, "Header: \"Test case review\"", 0.50, y + 2.92, 6.2, 0.24,
    size=SZ_SMALL, bold=True, color=NAVY)
txt(s, "Question: \"LBVOICESER-XXXX — <summary> (4 test cases). Proceed with script generation?\"",
    0.50, y + 3.15, 6.2, 0.28, size=SZ_SMALL, color=DARK)

for i, (label, clr) in enumerate([("Approve test cases", GREEN), ("Reject test cases", ORANGE)]):
    add_rect(s, 0.35, y + 3.60 + i * 0.55, 6.5, 0.46, fill=clr)
    txt(s, label, 0.50, y + 3.65 + i * 0.55, 6.2, 0.36,
        size=SZ_BODY, bold=True, color=WHITE)

# Right column: what happens on approve
sub(s, "On Approve — automatic actions (no further prompts):", 7.10, y + 0.05)
auto_steps = [
    ("1", NAVY,  "createJiraIssue",         "Type: Test · summary: '<parent> - Test case' · assignee: Bhimrao"),
    ("2", CYAN,  "transitionJiraIssue",      "Moves new Test issue → In Progress"),
    ("3", NAVY,  "editJiraIssue",            "Adds ticket to current active sprint"),
    ("4", GREEN, "Xray GraphQL (PowerShell)","addTestStep mutation for each test case — action/data/result"),
    ("5", CYAN,  "Inline confirm",           "✔ Created <key> → In Progress, Sprint 11, N steps added"),
    ("6", NAVY,  "Proceeds to Gate 2",       "Writes the Java class, no pause"),
]
for i, (num, col, title, desc) in enumerate(auto_steps):
    ty = y + 0.52 + i * 0.76
    pill(s, num, 7.10, ty + 0.10, 0.42, 0.42, bg=col)
    add_rect(s, 7.60, ty, 5.68, 0.68, fill=WHITE, line=MGREY)
    txt(s, title, 7.72, ty + 0.04, 5.44, 0.28, size=SZ_SMALL, bold=True, color=NAVY)
    txt(s, desc,  7.72, ty + 0.32, 5.44, 0.30, size=SZ_SMALL, italic=True, color=DGREY)

sub(s, "On Reject:", 7.10, y + 5.12)
txt(s, "No file written. Ticket re-enters the queue on the next sync.",
    7.10, y + 5.52, 6.0, 0.36, size=SZ_BODY, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Gate 2: Script Review
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Gate 2 — Script Review  (Human Approval)",
           "Claude writes the Java test class and shows it before any test runs.")

sub(s, "What Claude writes:", 0.35, y + 0.05)

# Code snippet
add_rect(s, 0.35, y + 0.50, 7.00, 3.80, fill=RGBColor(0x0D, 0x1A, 0x2E), line=CYAN)
code = [
    "@Feature(\"TTS POST API\")",
    "@XrayTest(key = \"LBVOICESER-1340\")",
    "@Requirement({\"LBVOICESER-1330\"})",
    "class TtsPostApiTest {",
    "",
    "  @Test",
    "  @DisplayName(\"No auth header returns 400\")",
    "  void missingAuthHeaderReturns400() {",
    "    TtsClient.synthesizeWithoutAuth(",
    "      TtsTestData.random())",
    "      .then().statusCode(400);",
    "  }",
    "  // ... 3 more @Test methods",
    "}",
]
for i, line in enumerate(code):
    txt(s, line, 0.50, y + 0.58 + i * 0.25, 6.7, 0.26,
        size=SZ_SMALL, color=GREEN, mono=True)

# Right: review gate
sub(s, "AskUserQuestion (UI dropdown):", 7.60, y + 0.05)
add_rect(s, 7.60, y + 0.50, 5.38, 0.62, fill=RGBColor(0xEB, 0xF3, 0xFB), line=CYAN)
txt(s, "Header: \"Script review\"", 7.75, y + 0.57, 5.1, 0.24, size=SZ_SMALL, bold=True, color=NAVY)
txt(s, "Question: \"Review src/.../XxxApiTest.java. Queue for combined run or discard?\"",
    7.75, y + 0.80, 5.1, 0.28, size=SZ_SMALL, color=DARK)

for i, (label, clr) in enumerate([("Approve — queue for combined run", GREEN),
                                   ("Reject — discard the file", ORANGE)]):
    add_rect(s, 7.60, y + 1.28 + i * 0.55, 5.38, 0.46, fill=clr)
    txt(s, label, 7.75, y + 1.33 + i * 0.55, 5.1, 0.36,
        size=SZ_BODY, bold=True, color=WHITE)

sub(s, "On Approve:", 7.60, y + 2.52)
bullets(s, [
    "Class name added to the combined run queue",
    "Maven does NOT run yet — waits for all tickets",
], 7.60, y + 2.96, 5.38, size=SZ_BODY, color=DARK, gap=0.46)

sub(s, "On Reject:", 7.60, y + 3.96)
bullets(s, [
    "Generated file is deleted from disk",
    "Jira comment posted: 'scripts rejected — re-evaluated next sync'",
    "No label stamped — ticket re-enters queue next sync",
], 7.60, y + 4.40, 5.38, size=SZ_SMALL, color=DARK, gap=0.40)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Gate 3: Combined Maven Run
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Gate 3 — Combined Maven Run  (Automatic)",
           "After ALL tickets are processed — one single mvn test command for all approved classes.")

sub(s, "Command executed (PowerShell — foreground, 5-minute timeout):", 0.35, y + 0.05)
add_rect(s, 0.35, y + 0.52, 12.6, 0.60, fill=RGBColor(0x0D, 0x1A, 0x2E), line=CYAN)
txt(s, "mvn test  -Dtest=TtsPostApiTest,TtsStreamPostApiTest,VoiceConfigApiTest",
    0.55, y + 0.64, 12.2, 0.40, size=SZ_BODY, color=GREEN, mono=True)

sub(s, "Rules:", 0.35, y + 1.30)
bullets(s, [
    "Never run per-ticket — one combined execution after ALL tickets are processed",
    "Always foreground (not background) — full output captured in the tool result",
    "PowerShell timeout: 300 000 ms  (5 minutes)",
    "pom.xml exec plugin opens the Allure report in a new window automatically",
], 0.35, y + 1.75, 12.6, size=SZ_BODY, color=DARK, gap=0.46)

sub(s, "Output captured — per-class result lines:", 0.35, y + 3.78)
add_rect(s, 0.35, y + 4.22, 12.6, 1.00, fill=RGBColor(0x0D, 0x1A, 0x2E), line=CYAN)
output_lines = [
    "[INFO] Tests run: 4, Failures: 0, Errors: 0  -- in com.laerdal.api.tests.TtsPostApiTest",
    "[ERROR] Tests run: 4, Failures: 2, Errors: 0 -- in com.laerdal.api.tests.TtsStreamPostApiTest",
    "[ERROR]   TtsStreamPostApiTest.missingVoiceFieldReturns200 -- Expected 200 but was 400",
]
for i, line in enumerate(output_lines):
    c = GREEN if "Failures: 0" in line else (ORANGE if "ERROR" in line else WHITE)
    txt(s, line, 0.55, y + 4.30 + i * 0.30, 12.2, 0.30, size=SZ_SMALL, color=c, mono=True)

sub(s, "Produces one Allure report covering all approved test classes.", 0.35, y + 5.40)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Post-Run Jira Updates
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Post-Run — Automatic Jira Updates  (Per Approved Class)",
           "Same turn, no user input needed. All label + transition calls run in parallel.")

sub(s, "Step 1 — always (every approved class):", 0.35, y + 0.05)
add_rect(s, 0.35, y + 0.52, 12.6, 0.56, fill=NAVY)
txt(s, "editJiraIssue  →  { \"labels\": [\"qa-auto-generated\"] }",
    0.55, y + 0.62, 12.2, 0.38, size=SZ_BODY, color=CYAN, mono=True)

# Two paths
sub(s, "IF  Failures + Errors == 0  (all tests pass):", 0.35, y + 1.28, color=GREEN)
add_rect(s, 0.35, y + 1.72, 5.9, 0.56, fill=GREEN)
txt(s, "transitionJiraIssue  →  transition ID 31  (Done)",
    0.55, y + 1.82, 5.6, 0.38, size=SZ_BODY, bold=True, color=WHITE, mono=True)
bullets(s, [
    "Ticket moves to Done automatically",
    "No human step required",
    "Jira Test issue (LBVOICESER-XXXX) also stays In Progress",
    "  until Xray import updates it",
], 0.35, y + 2.40, 5.9, size=SZ_SMALL, color=DARK, gap=0.40)

sub(s, "IF  Failures + Errors > 0  (any test fails):", 7.10, y + 1.28, color=ORANGE)
add_rect(s, 7.10, y + 1.72, 5.85, 0.56, fill=ORANGE)
txt(s, "addCommentToJiraIssue  →  ADF format",
    7.28, y + 1.82, 5.5, 0.38, size=SZ_BODY, bold=True, color=WHITE, mono=True)
bullets(s, [
    "Tags @reporter (account mention)",
    "Tests run: N | Failures: F | Errors: E",
    "Lists each failing test method name",
    "States expected vs actual HTTP status",
    "Requests: 'review Allure report and correct AC'",
    "Closes with: QA Automation",
], 7.10, y + 2.40, 5.85, size=SZ_SMALL, color=DARK, gap=0.40)

# Divider
add_rect(s, 6.87, y + 1.28, 0.04, 4.50, fill=MGREY)

sub(s, "Ticket stays in 'Ready for testing' — no auto-transition on failure.",
    7.10, y + 4.95, color=DGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — End-to-end flow summary
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "End-to-End Flow Summary",
           "From  start  command to Jira updated — everything in one Claude Code session.")

flow = [
    (NAVY,   "start",        "User types start in Claude Code chat"),
    (CYAN,   "JQL Fetch",    "All 'Ready for testing' tickets without processed labels"),
    (NAVY,   "Ctx Check",    "SUFFICIENT: has summary + desc + AC + method + endpoint + example"),
    (ORANGE, "INSUFFICIENT", "ADF comment → qa-context-requested label → excluded from next sync"),
    (GREEN,  "Gate 1",       "Test case table shown → Approve / Reject (AskUserQuestion)"),
    (CYAN,   "Auto-create",  "Jira Test issue → In Progress → Sprint → Xray steps"),
    (NAVY,   "Gate 2",       "Java class written → shown → Approve / Reject (AskUserQuestion)"),
    (GREEN,  "Gate 3",       "One mvn test covering all approved classes · Allure report"),
    (NAVY,   "Post-run",     "qa-auto-generated label · Done transition OR failure comment @reporter"),
]

COLS   = 5
ROWS   = 2
AVAIL_W = 12.70
AVAIL_H = 7.50 - y - 0.20
ARROW_W = 0.30
BOX_W  = (AVAIL_W - ARROW_W * (COLS - 1)) / COLS      # ~2.34
DESC_H = 1.10
BOX_H  = (AVAIL_H - DESC_H * ROWS) / ROWS             # box height per row

ROW_TOTAL = BOX_H + DESC_H                             # full row height incl desc

for i, (col, phase, desc) in enumerate(flow):
    c = i % COLS
    r = i // COLS
    x  = 0.32 + c * (BOX_W + ARROW_W)
    ty = y + 0.20 + r * ROW_TOTAL

    # coloured box
    add_rect(s, x, ty, BOX_W, BOX_H, fill=col)
    txt(s, phase, x, ty, BOX_W, BOX_H,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # description below the box
    txt(s, desc, x, ty + BOX_H + 0.06, BOX_W, DESC_H - 0.10,
        size=14, color=DGREY, align=PP_ALIGN.CENTER)

    # arrow (→) between boxes in same row — not after last in row
    if c < COLS - 1:
        txt(s, "→", x + BOX_W, ty + BOX_H / 2 - 0.18, ARROW_W, 0.38,
            size=16, bold=True, color=MGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Benefits
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Benefits",
           "Impact on QA speed, consistency, and developer feedback loop.")

benefits = [
    (GREEN,  "Zero manual Jira scanning",
     "Tickets in 'Ready for testing' are discovered automatically on every start."),
    (NAVY,   "Test cases in seconds",
     "From AC to a reviewed Java class in one session — no handoff delay."),
    (CYAN,   "One-time processing guarantee",
     "Label filter ensures each ticket is processed exactly once across all syncs."),
    (GREEN,  "Xray traceability built-in",
     "Jira Test issue + Xray steps created automatically on Gate 1 approve."),
    (NAVY,   "Instant failure feedback",
     "@reporter mentioned in Jira with exact failing method names within seconds."),
    (ORANGE, "Human control at every gate",
     "Gate 1 (plan) + Gate 2 (script) — no test runs without QA sign-off."),
]

BW = 5.95; BH = 1.50
for i, (col, title, desc) in enumerate(benefits):
    cx = 0.38 + (i % 2) * 6.50
    cy = y + 0.25 + (i // 2) * (BH + 0.15)
    add_rect(s, cx, cy, BW, BH, fill=WHITE, line=MGREY)
    add_rect(s, cx, cy, BW, 0.08, fill=col)
    txt(s, title, cx + 0.18, cy + 0.14, BW - 0.3, 0.44,
        size=SZ_SUB, bold=True, color=NAVY)
    txt(s, desc, cx + 0.18, cy + 0.62, BW - 0.3, 0.72,
        size=SZ_SMALL, italic=True, color=DGREY)



# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Future Scope: Fully Hands-Free CI/CD Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
y = header(s, "Future Scope  —  Fully Hands-Free CI/CD Pipeline",
           "Moving from a VS Code trigger to a fully automated GitHub Actions + Microsoft Teams workflow")

# ── Five pipeline stage boxes ─────────────────────────────────────────────────
STAGE_Y = y + 0.20
STAGE_H = 2.10
STAGE_W = 2.30
GAP     = 0.22
START_X = 0.32

stages = [
    (NAVY,  WHITE, "Developer",
     ["Moves Jira ticket", "→  'Ready for QA'", "(no manual trigger)"]),
    (CYAN,  NAVY,  "GitHub Actions",
     ["Jira webhook fires", "GHA workflow starts", "automatically"]),
    (ORANGE,WHITE, "Teams: Plan",
     ["Adaptive card sent", "to QA channel", "Approve / Reject plan"]),
    (NAVY,  WHITE, "Teams: Script",
     ["Adaptive card sent", "with generated file", "Approve / Reject script"]),
    (CYAN,  NAVY,  "Teams: Results",
     ["mvn test runs in CI", "Allure report link", "published to channel"]),
]

for i, (bg_col, fg, title, lines) in enumerate(stages):
    x = START_X + i * (STAGE_W + GAP)
    # header bar
    add_rect(s, x, STAGE_Y, STAGE_W, 0.52, fill=bg_col)
    txt(s, title, x, STAGE_Y, STAGE_W, 0.52,
        size=SZ_SUB, bold=True, color=fg, align=PP_ALIGN.CENTER)
    # body
    add_rect(s, x, STAGE_Y + 0.52, STAGE_W, STAGE_H - 0.52, fill=WHITE, line=MGREY)
    for j, line in enumerate(lines):
        txt(s, line, x + 0.10, STAGE_Y + 0.62 + j * 0.44, STAGE_W - 0.18, 0.40,
            size=SZ_BODY, color=DARK, align=PP_ALIGN.CENTER, italic=(j > 0))
    # arrow between stages
    if i < len(stages) - 1:
        ax = x + STAGE_W + 0.02
        txt(s, "▶", ax, STAGE_Y + STAGE_H / 2 - 0.18, GAP + 0.02, 0.38,
            size=18, bold=True, color=MGREY, align=PP_ALIGN.CENTER)



# ── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\INBDA2\Projects\API Automation suite\docs\QA-Automation-Demo-v7.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
